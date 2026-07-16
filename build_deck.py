#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build a skeleton MMN-focused lecture deck (Japanese main text)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

JP = "Yu Gothic"          # main JP font (falls back gracefully in PowerPoint)
JP_B = "Yu Gothic"        # bold uses same family
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xA4)
LIGHT = RGBColor(0xEF, 0xF3, 0xF8)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_font(run, size, bold=False, color=DARK, font=JP):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    # set East Asian font too
    rPr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def add_rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def add_text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        text, size, bold, color = ln.get("t",""), ln.get("s",18), ln.get("b",False), ln.get("c",DARK)
        p.alignment = ln.get("a", PP_ALIGN.LEFT)
        p.space_after = Pt(ln.get("sa", 4))
        p.space_before = Pt(ln.get("sb", 0))
        if ln.get("level"): p.level = ln["level"]
        r = p.add_run(); r.text = text
        _set_font(r, size, bold, color, ln.get("f", JP))
    return tb


def content_slide(title, kicker, bullets, footer=None, section_no=None):
    s = prs.slides.add_slide(BLANK)
    # top title band
    add_rect(s, 0, 0, SW, Inches(1.15), NAVY)
    add_rect(s, 0, Inches(1.15), SW, Pt(4), ACCENT)
    # kicker (section label)
    if kicker:
        add_text(s, Inches(0.5), Inches(0.12), Inches(11), Inches(0.3),
                 [{"t": kicker, "s": 12, "b": True, "c": RGBColor(0xBB,0xD3,0xEA)}])
    add_text(s, Inches(0.5), Inches(0.36), Inches(12.3), Inches(0.75),
             [{"t": title, "s": 26, "b": True, "c": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
    # body
    tb = s.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(12.1), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        lvl = b.get("level", 0)
        p.level = lvl
        p.space_after = Pt(b.get("sa", 7))
        p.space_before = Pt(b.get("sb", 2))
        marker = b.get("marker", "●" if lvl == 0 else "–")
        color = b.get("c", DARK if lvl == 0 else GREY)
        size = b.get("s", 18 if lvl == 0 else 15)
        bold = b.get("b", False)
        if marker:
            r0 = p.add_run(); r0.text = marker + " "
            _set_font(r0, size, True, b.get("mc", ACCENT if lvl==0 else GREY))
        r = p.add_run(); r.text = b["t"]
        _set_font(r, size, bold, color)
    # footer / source
    if footer:
        add_text(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4),
                 [{"t": footer, "s": 10.5, "c": GREY}])
    if section_no:
        add_text(s, Inches(12.5), Inches(0.12), Inches(0.7), Inches(0.3),
                 [{"t": section_no, "s": 12, "b": True, "c": RGBColor(0xBB,0xD3,0xEA), "a": PP_ALIGN.RIGHT}])
    return s


def section_divider(no, title, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, Inches(0.9), Inches(3.0), Inches(1.6), Pt(6), ACCENT)
    add_text(s, Inches(0.9), Inches(2.0), Inches(3), Inches(1.0),
             [{"t": no, "s": 80, "b": True, "c": RGBColor(0x3E,0x6A,0x99)}])
    add_text(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.3),
             [{"t": title, "s": 36, "b": True, "c": WHITE}])
    if subtitle:
        add_text(s, Inches(0.95), Inches(4.4), Inches(11.5), Inches(1.0),
                 [{"t": subtitle, "s": 18, "c": RGBColor(0xBB,0xD3,0xEA)}])
    return s


# ---------------- Slide 1: Title ----------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(4.55), SW, Pt(5), ACCENT)
add_text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.9),
         [{"t": "新ガイドラインを踏まえた", "s": 34, "b": True, "c": WHITE, "sa": 6},
          {"t": "多巣性運動ニューロパチーのアップデート", "s": 40, "b": True, "c": WHITE}])
add_text(s, Inches(0.95), Inches(4.75), Inches(11.5), Inches(1.4),
         [{"t": "天理よろづ相談所病院 神経筋疾患センター・脳神経内科", "s": 17, "c": RGBColor(0xCF,0xDD,0xEC), "sa": 4},
          {"t": "野寺 裕之", "s": 22, "b": True, "c": WHITE, "sa": 4},
          {"t": "武田薬品工業 全国研究会", "s": 14, "c": RGBColor(0xAF,0xC5,0xDD)}])
add_text(s, Inches(0.95), Inches(6.6), Inches(11), Inches(0.4),
         [{"t": "※ 骨子案（figure/データは要挿入・出典明記、国内未承認薬は適応外注記）", "s": 11, "c": RGBColor(0x8F,0xA8,0xC2)}])

# ---------------- Slide 2: COI ----------------
content_slide("COI の開示", "DISCLOSURE",
    [{"t": "演者氏名： 野寺 裕之", "s": 20, "b": True},
     {"t": "天理よろづ相談所病院 神経筋疾患センター長／脳神経内科 副部長", "s": 16, "level":1, "marker":""},
     {"t": "本講演に関連し、開示すべき COI 関係にある企業等の有無：（記入）", "s": 18, "sb":10},
     {"t": "＊本スライドは骨子。実際の COI 状態に合わせて記載してください。", "s": 13, "c":GREY, "marker":"", "sb":14}])

# ---------------- Slide 3: Agenda ----------------
content_slide("本日の内容", "AGENDA",
    [{"t": "MMN とは — 疾患概要・疫学・病態のアップデート", "b":True},
     {"t": "日本全国調査2024、補体（C2依存性）が病態の中心へ", "level":1, "marker":"–"},
     {"t": "診断のアップデートと「新ガイドライン」", "b":True, "sb":6},
     {"t": "診断基準の再検討（CBを超えて）、EAN/PNS 新ガイドライン近刊", "level":1, "marker":"–"},
     {"t": "治療のアップデート —「本命」補体阻害薬", "b":True, "sb":6},
     {"t": "empasiprubart（ARDA/EMPASSION）ほか、新規IgG療法への接続", "level":1, "marker":"–"}])

# ================= SECTION 1 =================
section_divider("1", "MMN とは", "疾患概要・疫学・病態のアップデート")

content_slide("疾患概念 — 純粋運動性の免疫介在性ニューロパチー", "1｜疾患概要",
    [{"t": "感覚障害を伴わない、左右非対称・上肢遠位優位の筋力低下と筋萎縮", "b":True},
     {"t": "緩徐〜階段状進行。運動神経伝導ブロック（CB）が特徴", "level":1, "marker":"–"},
     {"t": "約半数で抗GM1 IgM抗体陽性、IVIgが有効", "level":1, "marker":"–"},
     {"t": "cold paresis（寒冷麻痺）が83%（CIDP/PMAの4〜6倍）— 補助的特徴", "level":1, "marker":"–"},
     {"t": "鑑別の要は ALS/運動ニューロン疾患、および multifocal CIDP（MADSAM/LSS）", "b":True, "sb":8},
     {"t": "ALSとの臨床鑑別：筋萎縮に不釣合いな筋力低下、球麻痺なし、UMN徴候なし", "level":1, "marker":"–"}],
    footer="出典：Claytor et al. Muscle Nerve 2025; 71:512-534／日本CIDP/MMN診療ガイドライン2024 等より作図", section_no="1")

content_slide("疫学 — 日本全国調査2024（正式論文化）", "1｜疫学",
    [{"t": "全国推定患者数 約500名／有病率 0.4（人口10万対）", "b":True, "c":NAVY},
     {"t": "男女比 2:1、平均発症年齢 42歳（男性が若年発症）、上肢遠位に初発", "level":1, "marker":"–"},
     {"t": "抗GM1 IgM陽性率 約54%（英国22.7%・蘭/墺43%より高い）", "level":1, "marker":"–"},
     {"t": "海外の有病率は 0.3〜0.7/10万（希少疾患）", "sb":8},
     {"t": "HLA-DRB1*15 が高頻度（遺伝的背景）", "level":1, "marker":"–"}],
    footer="出典：Aotsuka Y, et al. Multifocal motor neuropathy in Japan. Muscle Nerve 2024;70:1027-1033 等より作図", section_no="1")

content_slide("病態① — nodopathy / paranodopathy", "1｜病態",
    [{"t": "CBの本態は「局所脱髄」から「ランビエ絞輪の機能障害」へ（Uncini）", "b":True},
     {"t": "paranodeミエリン剝離・絞輪伸長・Navチャネル障害・イオン恒常性破綻", "level":1, "marker":"–"},
     {"t": "神経興奮性研究：CB遠位で過分極、活動依存性CB", "level":1, "marker":"–"},
     {"t": "cold paresis は Na⁺/K⁺ポンプの温度感受性で説明", "level":1, "marker":"–"},
     {"t": "→ 「脱髄」だけでなく軸索・イオンチャネル機能障害という枠組み", "b":True, "c":ACCENT, "sb":8, "marker":"▶"}],
    footer="出典：Uncini A, Kuwabara S. J Neurol Neurosurg Psychiatry 2015;86:1186-95 等より作図（既存デッキの図を流用可）", section_no="1")

content_slide("病態② — 抗GM1 IgM と補体（C2依存性）＝治療標的の伏線", "1｜病態",
    [{"t": "抗GM1 IgM が古典経路（classical pathway）のみで補体を活性化", "b":True},
     {"t": "iPS由来運動ニューロンで補体依存性にCa恒常性破綻→軸索障害", "level":1, "marker":"–"},
     {"t": "【添付論文】MMN/抗MAG 51例のIgMによる補体活性化は全例 C2依存性", "b":True, "c":NAVY, "sb":8},
     {"t": "抗C2抗体で用量依存的に阻害、「C2バイパス」なし", "level":1, "marker":"–"},
     {"t": "→ C2 が有望な治療標的（＝empasiprubart の科学的基盤）", "level":1, "marker":"▶", "mc":GREEN, "c":GREEN, "b":True}],
    footer="出典：Budding K, et al. Complement activation ... depends on C2. Eur J Neurol 2025;32:e16541", section_no="1")

content_slide("病態③ — 補体活性化は臨床像と相関、阻害薬の前臨床根拠", "1｜病態（新）",
    [{"t": "MMN 137例の血清 × iPS由来運動ニューロンモデル", "b":True},
     {"t": "IgM結合↑→C3沈着↑が、筋力低下(MRC)・必要IVIg用量・軸索障害の程度と相関", "level":1, "marker":"–"},
     {"t": "補体活性化は振動覚異常・腕神経叢MRI異常とも相関", "level":1, "marker":"–"},
     {"t": "IVIgは補体活性化を34–54%しか抑制しない", "b":True, "c":NAVY, "sb":8},
     {"t": "特異的補体阻害薬はモデルで89.1–98.7%抑制 → 補体阻害の前臨床的根拠", "level":1, "marker":"▶", "mc":GREEN, "c":GREEN, "b":True}],
    footer="出典：Krijgsman D, … van der Pol WL. Neurol Neuroimmunol Neuroinflamm 2026;13(1):e200482", section_no="1")

content_slide("病態④ — 血液神経関門(BNB)と軸索膜（本邦グループの知見）", "1｜病態（新）",
    [{"t": "患者IgGによる血液神経関門(BNB)破綻：CIDPと分子機序が異なる", "b":True},
     {"t": "MMNでは内皮のTNF-α・VCAM-1上昇が主体（典型CIDPはCCL20/ICAM-1、多巣性CIDPはGM-CSF等）", "level":1, "marker":"–"},
     {"t": "軸索膜の不安定性：fast K⁺電流低下・superexcitability", "b":True, "c":NAVY, "sb":8},
     {"t": "paranodeミエリンの緩みを示唆、不随意指運動例でより顕著", "level":1, "marker":"–"},
     {"t": "→ MMNは「CB＋軸索」だけでなく、より広汎な有髄線維障害のスペクトラム", "level":1, "marker":"▶", "c":ACCENT, "b":True, "sb":6}],
    footer="出典：Shimizu F, et al. Int J Mol Sci 2026;27(2):1088（本邦）／Krarup C, et al. Clin Neurophysiol 2025;173:229-238", section_no="1")

# ================= SECTION 2 =================
section_divider("2", "診断のアップデート", "診断基準の再検討と EAN/PNS 新ガイドライン")

content_slide("診断プロセスと現行 EFNS/PNS 2010 基準", "2｜診断",
    [{"t": "臨床：2神経以上・1か月以上の非対称性運動麻痺、客観的感覚障害なし", "b":True},
     {"t": "支持：上肢優位・腱反射低下・脳神経正常・筋痙攣/線維束収縮・免疫療法反応", "level":1, "marker":"–"},
     {"t": "電気診断：絞扼部以外での definite/probable CB", "level":1, "marker":"–"},
     {"t": "補助：抗GM1 IgM、末梢神経肥厚（エコー/MRI）、軽度CSF蛋白上昇", "level":1, "marker":"–"},
     {"t": "疑い例では IVIg 診断的治療で反応性を確認することもしばしば", "sb":8, "marker":"●"}],
    footer="出典：Joint Task Force EFNS/PNS. J Peripher Nerv Syst 2010;15:295-301（表を流用）", section_no="2")

content_slide("診断アップデート① — 基準の感度比較（Doneddu 2024）", "2｜診断",
    [{"t": "EFNS/PNS 2010 vs AANEM：MMN 53例 vs 対照280例", "b":True},
     {"t": "感度：EFNS/PNS 47%（definite）/57%（prob+def）＞ AANEM 28%/53%", "level":1, "marker":"–"},
     {"t": "支持基準を加えると EFNS/PNS の prob+def 感度は 64% に", "level":1, "marker":"–"},
     {"t": "特異度：AANEM 100%、EFNS/PNS 97〜98.5%（ともに高い）", "level":1, "marker":"–"},
     {"t": "→ EFNS/PNSは高感度・わずかに低特異度。NCSの拡張（多神経検索）を推奨", "b":True, "c":ACCENT, "marker":"▶", "sb":8}],
    footer="出典：Doneddu PE, et al. Eur J Neurol 2024;31:e16444／同 Italian database 2024;31:e16248", section_no="2")

content_slide("診断アップデート② — CBを超えた「拡張基準」（Doneddu 2025）", "2｜診断",
    [{"t": "MMN 70例 vs 対照359例", "b":True},
     {"t": "CBのない区間でも71%にCB以外の伝導異常（時間的分散47%・遠位CMAP持続延長31%・CV低下26% 等）", "level":1, "marker":"–"},
     {"t": "これらを取り込むと prob/def 感度が +26%（p<0.001）、特異度低下は最小", "b":True, "c":NAVY, "sb":8},
     {"t": "「CBが唯一の脱髄指標ではない」— 新ガイドライン改訂の論点に", "level":1, "marker":"▶", "mc":GREEN, "c":GREEN, "b":True},
     {"t": "関連：CB定義のばらつき／CB依存度への問題提起が相次ぐ（2025 レター）", "level":1, "marker":"–", "sb":6}],
    footer="出典：Doneddu PE, et al. Nerve Conduction Abnormalities Beyond CB in MMN. Eur J Neurol 2025;32:e70300", section_no="2")

content_slide("診断アップデート③ — 実データが示す電気診断の最適化（2025）", "2｜診断（新）",
    [{"t": "CB定義は「振幅30%減」が「面積30%減」より高感度（78.7% vs 63.8%）", "b":True},
     {"t": "CB(振幅)＋時間的分散(TD)＋F波延長/消失を独立指標にすると感度78.7→91.5%", "level":1, "marker":"▶", "mc":GREEN, "c":GREEN, "b":True},
     {"t": "TDは感度59.6%/特異度94.3%、F波異常は感度42.6%/特異度96.9%（英・韓 多施設）", "level":1, "marker":"–"},
     {"t": "予後：IVIg反応不良例は発症時から罹患神経数が多くTDが多い（CB単独では判別不可）", "b":True, "c":NAVY, "sb":8},
     {"t": "→ CBだけでなくTD・F波・罹患神経数を系統的に評価", "level":1, "marker":"–", "c":ACCENT}],
    footer="出典：Rajabally YA, et al. Eur J Neurol 2025;32:e70361／Posa A, et al. IBRO Neurosci Rep 2026;20:474-479", section_no="2")

content_slide("画像・バイオマーカーのアップデート", "2｜診断",
    [{"t": "神経エコー：非対称性の神経腫大を高感度に検出、ALS鑑別に有用", "b":True},
     {"t": "低コスト・非侵襲。narrative reviewは「可能なら第一選択の画像」と提言", "level":1, "marker":"–"},
     {"t": "MR neurography：神経腫大は MADSAM で顕著、MMNで乏しい→鑑別の手がかり", "sb":6, "marker":"●", "b":True},
     {"t": "バイオマーカー：血清NfL は ALS＜ 鑑別補助、peripherin（軸索障害）も注目", "sb":6, "marker":"●", "b":True},
     {"t": "現行 EFNS/PNS はエコー記載なし → 新GLでの正式採用が期待される", "level":1, "marker":"–", "c":ACCENT}],
    footer="出典：Pitman J, et al. Muscle Nerve 2024;71:293-308／Karam C. Muscle Nerve 2025;73:86-92 等", section_no="2")

content_slide("神経エコー — MMN・純粋運動CIDP・典型CIDPは連続スペクトラム", "2｜診断（新）",
    [{"t": "MMN19／motor CIDP15／典型CIDP117 を神経エコー＋NCSで比較", "b":True},
     {"t": "神経断面積(CSA)：典型CIDP＞MPred-CIDP＞PM-CIDP＞MMN の順で小さくなる", "level":1, "marker":"–"},
     {"t": "伝導速度は逆順（MMNで最も速い）→ 3疾患は連続スペクトラム", "level":1, "marker":"–"},
     {"t": "motor CIDPは腕神経叢が近位優位に腫大（MMN・典型CIDPは均等）", "level":1, "marker":"–"},
     {"t": "→ エコー＋NCS＋治療反応で疾患スペクトラム上の位置づけが可能", "level":1, "marker":"▶", "c":ACCENT, "b":True, "sb":6}],
    footer="出典：Niu J, et al. Nerve Ultrasound of MMN and Motor/Typical CIDP. Can J Neurol Sci 2025", section_no="2")

content_slide("鑑別診断 — 手初発ALS と 運動優位ニューロパチー", "2｜診断（新）",
    [{"t": "MMN vs 手初発ALS（Fang 2026）", "b":True, "c":NAVY},
     {"t": "MMNは若年発症(43 vs 59歳)・低CK・高IgM。CBはMMN 87.5% vs 手初発ALS 31%", "level":1, "marker":"–"},
     {"t": "手初発ALSの約1/3でCK上昇、MMNは全例CK正常", "level":1, "marker":"–"},
     {"t": "LMND様重症運動ニューロパチーで新規セプチン多量体自己抗体（Arlt 2026, Brain）", "b":True, "c":NAVY, "sb":8},
     {"t": "MMN検証群は陰性だが、運動優位ニューロパチー×LMND鑑別で今後注目", "level":1, "marker":"–"}],
    footer="出典：Fang SY, et al. J Chin Med Assoc 2026;89:437-444／Arlt FA, et al. Brain 2026;awag183", section_no="2")

content_slide("診断遅延・誤診の実態 —「気づかれない」MMN", "2｜診断",
    [{"t": "英国多施設調査：初診からの診断遅延が50%超で1年以上", "b":True},
     {"t": "argenxグローバル医師調査（250名/9か国）", "sb":8, "b":True, "c":NAVY},
     {"t": "約70%が別の末梢神経疾患と最初に誤診（ALS誤診50%、CIDP51%）", "level":1, "marker":"–"},
     {"t": "31%はどのガイドラインも参照せず、1/3がステロイドを1st-line（非推奨）", "level":1, "marker":"–"},
     {"t": "診断遅延：医師推計 約1.75年／患者報告 約6年（両者に乖離、両方提示）", "level":1, "marker":"–", "c":GREY}],
    footer="出典：Rajabally YA, et al. JPNS 2025;30:e70018／argenx PNS 2025 ポスター（医師・患者調査）", section_no="2")

content_slide("新ガイドライン — EAN/PNS MMN ガイドライン（近刊）", "2｜新ガイドライン",
    [{"t": "現行標準は依然 2010年 EFNS/PNS 第1改訂版", "b":True},
     {"t": "改訂版：Task Force chair S. Goedee（Utrecht）", "b":True, "c":NAVY, "sb":8},
     {"t": "EAN公式は「Guideline Proposal accepted（初期段階）」、2025–2026発刊見込み", "level":1, "marker":"–"},
     {"t": "＊正式な発刊日・内容は未公開 → 本項は「予想」として提示", "level":1, "marker":"–", "c":GREY},
     {"t": "予想される変更点（未確定）", "b":True, "sb":8},
     {"t": "CB以外の伝導異常の位置づけ／神経エコーの支持基準採用／SCIg・補体阻害への言及", "level":1, "marker":"–"}],
    footer="出典：EAN Ongoing Guidelines（ean.org）。発表直前に発刊状況の再確認を推奨", section_no="2")

content_slide("参考：2021 EAN/PNS CIDPガイドライン と 日本GL2024", "2｜新ガイドライン",
    [{"t": "2021 CIDPガイドライン：「atypical CIDP」→「CIDP variants」に整理", "b":True},
     {"t": "Lewis-Sumner/MADSAM は「multifocal CIDP」と改称 → MMNとの厳密な鑑別が重要", "level":1, "marker":"–"},
     {"t": "日本GL2024（要点）", "b":True, "c":NAVY, "sb":8},
     {"t": "IVIg導入(1C)・維持(1B)を強く推奨、SCIg維持は条件付き(2C・国内保険適用外)", "level":1, "marker":"–"},
     {"t": "免疫抑制療法は行わないことを強く推奨(1B)、ステロイドは増悪例あり非推奨", "level":1, "marker":"–"},
     {"t": "補体阻害薬・FcRnは「今後の臨床試験に期待」→ 次章が最新の続報", "level":1, "marker":"▶", "c":ACCENT}],
    footer="出典：Van den Bergh PYK, et al. EAN/PNS CIDP guideline 2021／日本CIDP/MMN診療ガイドライン2024", section_no="2")

# ================= SECTION 3 =================
section_divider("3", "治療のアップデート", "「本命」補体阻害薬と新規IgG療法")

content_slide("治療の現状と限界 — IVIg / SCIg", "3｜治療",
    [{"t": "IVIgは唯一の承認治療・第一選択（実臨床の奏効率は高い：英国97.8%）", "b":True},
     {"t": "SCIgはIVIgと同等の有効性・良好な忍容性（国内は未承認）", "level":1, "marker":"–"},
     {"t": "限界：維持療法下でも軸索変性が緩徐に進行、経時的に効果減弱", "b":True, "c":NAVY, "sb":8},
     {"t": "約90%の患者が維持IVIg下でも軸索変性が進行（患者調査）", "level":1, "marker":"–"},
     {"t": "→ 病態を修飾する新規治療（補体阻害）への強い期待", "level":1, "marker":"▶", "mc":GREEN, "c":GREEN, "b":True, "sb":6}],
    footer="出典：Rajabally YA, et al. JPNS 2025;30:e70018／Claytor et al. Muscle Nerve 2025 等", section_no="3")

content_slide("IVIg の最適化と長期アウトカム（2025–2026）", "3｜治療（新）",
    [{"t": "高頻度IVIg（年2回超）で遠位CMAP振幅・mRSが有意改善、CB数も減少", "b":True},
     {"t": "過少治療だった症例でも改善余地（B=2.88, P<0.001／中国・地方3施設）", "level":1, "marker":"–"},
     {"t": "個別最適化用量で MMN-RODS が90.6%(29/32)で6年以上改善維持", "b":True, "c":NAVY, "sb":8},
     {"t": "従来の「経時的に低下」像と対照的。用量は施設間差大（18.8〜33.9 g/週）", "level":1, "marker":"–"},
     {"t": "→ 早期・十分量・個別最適化が鍵（減量判断のLLM支援も報告）", "level":1, "marker":"▶", "c":ACCENT, "b":True, "sb":6}],
    footer="出典：Zhou L, et al. BMC Neurol 2026／Noushad MA, et al. JPNS 2025;30:e70079／Burton E, et al. Transfus Med 2026", section_no="3")

content_slide("残された課題 — 非運動症状の負担 と SCIg の実際", "3｜治療（新）",
    [{"t": "6年追跡：運動障害(INCAT)は安定だが、疲労17→58%・うつ8→25%に増悪", "b":True, "c":NAVY},
     {"t": "免疫治療は運動を安定させるが非運動症状は防げない → 疲労・メンタルのスクリーニング", "level":1, "marker":"–"},
     {"t": "IVIg→SCIg 切替：MMNの77%が12か月継続・安定、満足度高い", "b":True, "sb":8},
     {"t": "ただし40%で増量など治療強化が必要（本邦は保険適用外）", "level":1, "marker":"–"},
     {"t": "→ 多職種での包括的マネジメントとSCIgという在宅選択肢", "level":1, "marker":"▶", "c":ACCENT, "b":True, "sb":6}],
    footer="出典：Basta I, et al. Front Hum Neurosci 2026;20:1824742／Gingele S, et al. Neurol Ther 2026;15:165-177", section_no="3")

content_slide("なぜ補体阻害か —「本命」の病態的根拠", "3｜治療",
    [{"t": "抗GM1 IgM →（古典経路・C2依存性）→ 補体活性化", "b":True},
     {"t": "ランビエ絞輪で膜侵襲複合体(MAC)形成 → Navチャネル障害・軸索障害", "level":1, "marker":"–"},
     {"t": "IVIgは補体沈着を減らすが不十分 → 上流での選択的補体阻害が合理的", "level":1, "marker":"–"},
     {"t": "開発の主戦場は C2（empasiprubart）と 活性化C1s（DNTH103）", "b":True, "c":NAVY, "sb":10, "marker":"▶"},
     {"t": "＊いずれも国内未承認・臨床試験段階（適応外）", "s":13, "c":GREY, "marker":"", "sb":8}],
    footer="出典：Budding K, et al. Eur J Neurol 2025;32:e16541 ほか", section_no="3")

content_slide("補体が本命、FcRn阻害はMMNには不向き", "3｜治療（新）",
    [{"t": "FcRn阻害（nipocalimab, efgartigimod, rozanolixizumab）はIgG自己抗体病で有効", "b":True},
     {"t": "重症筋無力症で確立、CIDPで開発中（efgartigimodは再発予防効果）", "level":1, "marker":"–"},
     {"t": "MMNではFcRn阻害は無効と考えられる ── MMNが主にIgM介在性の病態のため", "b":True, "c":NAVY, "sb":8, "marker":"●", "mc":NAVY},
     {"t": "MMNは「IgM × 補体」主導 → 標的は補体（C2 / 活性化C1s）が理にかなう", "level":1, "marker":"▶", "mc":GREEN, "c":GREEN, "b":True},
     {"t": "＊FcRn阻害薬・補体阻害薬とも国内はMMN未承認（適応外）", "s":13, "c":GREY, "marker":"", "sb":8}],
    footer="出典：Khateb M, Bril V. Nipocalimab and other FcRn blockers in neuromuscular disorders. Pharmacol Ther 2026;286:109071", section_no="3")

content_slide("empasiprubart（抗C2）— 第2相 ARDA 試験", "3｜治療（本命）",
    [{"t": "対象：IVIg依存の probable/definite MMN 計54名", "b":True},
     {"t": "27名×2コホート、各2:1（実薬n=18/プラセボn=9）、16週二重盲検", "level":1, "marker":"–"},
     {"t": "主要結果（EAN2025/PNS2025）", "b":True, "c":NAVY, "sb":8},
     {"t": "IVIg再投与リスク 約91%減：コホート1 HR 0.09（95%CI 0.02–0.44）", "level":1, "marker":"–", "c":GREEN, "b":True},
     {"t": "握力改善（コホート2：+19.89 vs +0.78 kPa）、血清NfL低下（軸索保護示唆）", "level":1, "marker":"–"},
     {"t": "安全性は概ね良好（多くは軽〜中等度）。MMNで過去最大の介入試験", "level":1, "marker":"–"}],
    footer="出典：ARDA 第2相（NCT05225675）argenx PNS 2025 ポスター。査読論文化は要確認・国内未承認", section_no="3")

content_slide("empasiprubart — 第3相 EMPASSION／他の補体阻害薬", "3｜治療（本命）",
    [{"t": "EMPASSION（NCT06742190）：empasiprubart vs IVIg、約154名、進行中", "b":True},
     {"t": "主要評価は握力変化（CT.gov）／MMN-RODS（企業資料）で出典間に食い違い→要確認", "level":1, "marker":"–", "c":GREY},
     {"t": "DNTH103（抗活性化C1s、皮下・2週毎）：第2相 MoMeNtum（MMN36名）、2H2026データ", "b":True, "sb":8},
     {"t": "riliprubart（抗C1s, Sanofi）はMMN非対象（CIDP専用・難治CIDP第3相は無益性中止）", "level":1, "marker":"–"},
     {"t": "eculizumab（抗C5）は小規模オープン試験で効果限定的（上流阻害へ主流移行）", "level":1, "marker":"–"}],
    footer="出典：ClinicalTrials.gov（NCT06742190/NCT06537999）／Sanofi press 2026／Fitzpatrick 2011", section_no="3")

# Complement inhibitor comparison table slide
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, Inches(1.15), NAVY)
add_rect(s, 0, Inches(1.15), SW, Pt(4), ACCENT)
add_text(s, Inches(0.5), Inches(0.12), Inches(11), Inches(0.3),
         [{"t": "3｜治療（本命）", "s": 12, "b": True, "c": RGBColor(0xBB,0xD3,0xEA)}])
add_text(s, Inches(0.5), Inches(0.36), Inches(12.3), Inches(0.75),
         [{"t": "補体阻害薬の一覧（MMN）", "s": 26, "b": True, "c": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
from pptx.util import Cm
rows, cols = 5, 6
tbl_x, tbl_y, tbl_w, tbl_h = Inches(0.5), Inches(1.5), Inches(12.33), Inches(4.4)
gtbl = s.shapes.add_table(rows, cols, tbl_x, tbl_y, tbl_w, tbl_h).table
headers = ["薬剤", "標的", "経路", "投与", "MMN試験", "段階/結果"]
data = [
    ["empasiprubart\n(ARGX-117)", "C2", "古典＋レクチン", "静注", "ARDA(P2)→\nEMPASSION(P3)", "IVIg再投与\n約91%減／P3進行中"],
    ["DNTH103", "活性化C1s", "古典のみ", "皮下・2週毎", "MoMeNtum(P2)", "募集中／2H2026"],
    ["riliprubart\n(SAR445088)", "C1s", "古典のみ", "静注", "なし（CIDP専用）", "難治CIDP第3相は中止"],
    ["eculizumab", "C5", "終末", "静注", "小規模オープン\n(2011)", "効果限定的"],
]
colw = [Inches(2.3), Inches(1.5), Inches(2.0), Inches(1.8), Inches(2.4), Inches(2.33)]
for j, w in enumerate(colw):
    gtbl.columns[j].width = w
for j, h in enumerate(headers):
    c = gtbl.cell(0, j); c.text = ""
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = h; _set_font(r, 14, True, WHITE)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
for i, row in enumerate(data, start=1):
    for j, val in enumerate(row):
        c = gtbl.cell(i, j); c.text = ""
        c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
        tf = c.text_frame; tf.word_wrap = True
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        first = True
        for part in val.split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = part
            hl = (j == 5 and i == 1)
            _set_font(r, 12.5, hl, GREEN if hl else DARK)
add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
         [{"t": "＊すべて国内未承認・臨床試験段階（適応外）。数値は学会発表主体で査読論文化は要確認。", "s": 12, "c": GREY},
          {"t": "出典：ClinicalTrials.gov／argenx・Sanofi 学会発表・press（2024–2026）", "s": 10.5, "c": GREY}])

# ---- New IgG therapies bridge (to existing corporate deck) ----
content_slide("新規 IgG 療法への接続", "3｜治療",
    [{"t": "IgG補充療法の選択肢：IVIg／従来SCIg／facilitated SCIg（fSCIG）", "b":True},
     {"t": "在宅自己注射・全身性副作用の軽減・通院負担の低減", "level":1, "marker":"–"},
     {"t": "（本邦の適応・製剤の位置づけは各製品情報に準拠）", "level":1, "marker":"–", "c":GREY},
     {"t": "＝ ここから既存デッキの企業パート（新規IgG製剤スライド）へ接続", "b":True, "c":ACCENT, "marker":"▶", "sb":10},
     {"t": "＊製品スライドは適応・注意事項・COIに従い挿入", "s":13, "c":GREY, "marker":"", "sb":8}],
    footer="※ 既存30枚デッキの slide 19–30（HyQvia 等）を後続に配置", section_no="3")

# ---- Take home ----
content_slide("まとめ — Take home messages", "SUMMARY",
    [{"t": "MMNは誤診・診断遅延が多い（ALS誤診50%）— 純粋運動・非対称・CBを想起", "b":True},
     {"t": "診断はCB至上主義から拡張へ：CB以外の伝導異常＋神経エコー＋NfLが補助", "b":True, "sb":6},
     {"t": "新ガイドライン（EAN/PNS）が近刊予定 — 診断基準・エコー・治療の更新に注目", "b":True, "sb":6},
     {"t": "病態の中心は補体（C2依存性）活性化", "b":True, "sb":6},
     {"t": "本命：補体阻害薬 empasiprubart がIVIg再投与を約91%減（第3相進行中）", "b":True, "c":GREEN, "mc":GREEN, "marker":"▶", "sb":6},
     {"t": "IVIg/SCIg（新規IgG療法）は引き続き治療の基盤", "b":True, "sb":6}],
    footer="国内未承認薬・臨床試験段階のデータを含みます（適応外）", section_no="")

# ---- References ----
content_slide("主要文献", "REFERENCES",
    [{"t":"Claytor B, Polston D, Li Y. MMN: A Narrative Review. Muscle Nerve 2025;71:512-534.", "s":13, "marker":"1."},
     {"t":"Rajabally YA, et al. Diagnosis and Management of MMN in the UK. JPNS 2025;30:e70018.", "s":13, "marker":"2.", "sb":2},
     {"t":"Budding K, et al. Complement activation by IgM ... depends on C2. Eur J Neurol 2025;32:e16541.", "s":13, "marker":"3.", "sb":2},
     {"t":"日本神経学会ほか. CIDP/MMN 診療ガイドライン2024.", "s":13, "marker":"4.", "sb":2},
     {"t":"Aotsuka Y, et al. MMN in Japan: nationwide survey. Muscle Nerve 2024;70:1027-1033.", "s":13, "marker":"5.", "sb":2},
     {"t":"Doneddu PE, et al. EFNS/PNS vs AAEM criteria. Eur J Neurol 2024;31:e16444.", "s":13, "marker":"6.", "sb":2},
     {"t":"Doneddu PE, et al. NCS Abnormalities Beyond CB in MMN. Eur J Neurol 2025;32:e70300.", "s":13, "marker":"7.", "sb":2},
     {"t":"Van den Bergh PYK, et al. EAN/PNS CIDP guideline (2nd revision) 2021.", "s":13, "marker":"8.", "sb":2},
     {"t":"ARDA(NCT05225675)／EMPASSION(NCT06742190)／MoMeNtum(NCT06537999)；EAN Ongoing Guidelines(ean.org).", "s":13, "marker":"9.", "sb":2}],
    footer=None, section_no="")

content_slide("主要文献（2025–2026 アップデート）", "REFERENCES",
    [{"t":"Rajabally YA, et al. Variably Defined CB, TD and Other Electrophysiological Abnormalities in MMN. Eur J Neurol 2025;32:e70361.", "s":12.5, "marker":"10."},
     {"t":"Posa A, et al. Patterns of motor nerve demyelination and prognostic significance in MMN. IBRO Neurosci Rep 2026;20:474-479.", "s":12.5, "marker":"11.", "sb":2},
     {"t":"Fang SY, et al. Differentiating MMN from hand-onset ALS. J Chin Med Assoc 2026;89:437-444.", "s":12.5, "marker":"12.", "sb":2},
     {"t":"Arlt FA, et al. Septin multimer autoantibodies in severe motor neuropathy. Brain 2026;awag183.", "s":12.5, "marker":"13.", "sb":2},
     {"t":"Niu J, et al. Nerve Ultrasound of MMN and Motor/Typical CIDP. Can J Neurol Sci 2025.", "s":12.5, "marker":"14.", "sb":2},
     {"t":"Krijgsman D, et al. IgM Anti-Ganglioside Binding and Complement Activation in iPSC-MN Model for MMN. Neurol Neuroimmunol Neuroinflamm 2026;13:e200482.", "s":12.5, "marker":"15.", "sb":2},
     {"t":"Shimizu F, et al. Blood-Nerve Barrier Breakdown Induced by IgG in CIDP and MMN. Int J Mol Sci 2026;27:1088.", "s":12.5, "marker":"16.", "sb":2},
     {"t":"Zhou L, et al. IVIg frequency in MMN: nerve conduction outcomes. BMC Neurol 2026. ／ Noushad MA, et al. Long-Term Outcomes via MMN-RODS. JPNS 2025;30:e70079.", "s":12.5, "marker":"17.", "sb":2},
     {"t":"Basta I, et al. Long-term motor stability but increasing non-motor burden in MMN. Front Hum Neurosci 2026;20:1824742. ／ Gingele S, et al. IV→SC Ig in CIDP and MMN. Neurol Ther 2026;15:165-177.", "s":12.5, "marker":"18.", "sb":2},
     {"t":"Khateb M, Bril V. Nipocalimab and other FcRn blockers in neuromuscular disorders. Pharmacol Ther 2026;286:109071.", "s":12.5, "marker":"19.", "sb":2}],
    footer=None, section_no="")

prs.save("MMN_update_2026_skeleton_ja.pptx")
print("Saved. Slides:", len(prs.slides._sldIdLst))
